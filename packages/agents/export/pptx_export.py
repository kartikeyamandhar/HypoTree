"""PowerPoint export for hypothesis trees and workplans."""
from __future__ import annotations

import logging
import os
from typing import Optional

logger = logging.getLogger(__name__)


def export_to_pptx(tree_data: dict, output_path: str) -> str:
    """Generate a PowerPoint deck from a hypothesis tree."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        logger.error("python-pptx not installed")
        raise

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    bg_color = RGBColor(0x0f, 0x11, 0x17)
    text_color = RGBColor(0xe8, 0xea, 0xf0)
    muted_color = RGBColor(0x9c, 0xa3, 0xaf)
    accent_color = RGBColor(0x63, 0x66, 0xf1)

    def set_slide_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def add_text(slide, left, top, width, height, text, font_size=12, color=text_color, bold=False):
        from pptx.util import Inches, Pt
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        return txBox

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, 1, 2, 11, 1, "HypoTree Analysis", 36, accent_color, True)
    add_text(slide, 1, 3.2, 11, 0.5, tree_data.get("question", ""), 20, text_color)
    add_text(slide, 1, 4.2, 11, 0.5,
        f"{tree_data.get('industry', '')} | {tree_data.get('company', '')} | "
        f"{tree_data.get('classification', {}).get('question_type', '').replace('_', ' ').title()}",
        14, muted_color)

    # Slide 2: Tree overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, 0.5, 0.3, 5, 0.5, "Hypothesis Tree Overview", 24, accent_color, True)

    root = tree_data.get("root", {})
    add_text(slide, 0.5, 1.0, 12, 0.5, f"Root: {root.get('statement', '')}", 14, text_color, True)

    y = 1.7
    for i, child in enumerate(root.get("children", [])[:6]):
        add_text(slide, 1.0, y, 11, 0.4, f"{i+1}. {child.get('statement', '')}", 12, text_color)
        y += 0.5

    # Slide 3: Key findings (stress test summary)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, 0.5, 0.3, 5, 0.5, "Red Team Findings", 24, RGBColor(0xef, 0x44, 0x44), True)

    report = tree_data.get("stress_test_report", {})
    add_text(slide, 0.5, 1.0, 12, 0.4,
        f"{report.get('critical_count', 0)} Critical | {report.get('warning_count', 0)} Warnings | {report.get('note_count', 0)} Notes",
        16, text_color)

    y = 1.8
    for critique in report.get("critiques", [])[:5]:
        sev = critique.get("severity", "note")
        color = RGBColor(0xfc, 0xa5, 0xa5) if sev == "critical" else RGBColor(0xfc, 0xd3, 0x4d) if sev == "warning" else muted_color
        add_text(slide, 0.5, y, 12, 0.4, f"[{sev.upper()}] {critique.get('claim_challenged', '')}", 11, color)
        y += 0.5

    # Slide 4: Workplan
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_text(slide, 0.5, 0.3, 5, 0.5, "Workplan", 24, RGBColor(0xf5, 0x9e, 0x0b), True)

    workplan = tree_data.get("workplan", {})
    add_text(slide, 0.5, 1.0, 12, 0.4,
        f"{len(workplan.get('workstreams', []))} Workstreams | {workplan.get('total_loe', 0):.0f}h Total | {workplan.get('estimated_weeks', 0)} Weeks",
        16, text_color)

    y = 1.8
    for ws in workplan.get("workstreams", [])[:6]:
        add_text(slide, 0.5, y, 12, 0.4,
            f"{ws.get('id', '')}: {ws.get('name', '')} ({ws.get('total_loe', 0):.0f}h, {len(ws.get('items', []))} analyses)",
            12, text_color)
        y += 0.45

    prs.save(output_path)
    logger.info("Exported PPTX to %s", output_path)
    return output_path
